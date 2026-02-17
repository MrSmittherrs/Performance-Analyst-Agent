"""
Slide Deck Generator: Creates a professional PowerPoint presentation
from YouTube trend analysis data.

Styled to match Inspired Testing corporate branding:
- White/light backgrounds with gold (#C2A269) accent
- Helvetica for headings, Calibri for body text
- Clean, minimal corporate layout

Usage:
    python tools/slide_generator.py
    python tools/slide_generator.py --data-dir .tmp
"""

import os
import sys
import json
import argparse
from datetime import datetime
from dotenv import load_dotenv

import matplotlib
matplotlib.use("Agg")  # Non-interactive backend
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ---------------------------------------------------------------------------
# Brand palette — Inspired Testing style
# ---------------------------------------------------------------------------

COLORS = {
    "bg_white": RGBColor(0xFF, 0xFF, 0xFF),
    "bg_light": RGBColor(0xF5, 0xF5, 0xF5),
    "bg_card": RGBColor(0xF0, 0xF0, 0xF0),
    "gold": RGBColor(0xC2, 0xA2, 0x69),          # Primary accent
    "charcoal": RGBColor(0x2D, 0x2E, 0x31),       # Sub-headings
    "text_dark": RGBColor(0x26, 0x26, 0x26),       # Body text
    "text_body": RGBColor(0x3B, 0x38, 0x38),       # Secondary body
    "text_muted": RGBColor(0x66, 0x66, 0x66),      # Page numbers, captions
    "black": RGBColor(0x00, 0x00, 0x00),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "success": RGBColor(0x2E, 0x7D, 0x32),         # Green — corporate
    "warning": RGBColor(0xC6, 0x28, 0x28),          # Red — corporate
    "table_header": RGBColor(0x2D, 0x2E, 0x31),     # Dark charcoal header
    "table_header_text": RGBColor(0xFF, 0xFF, 0xFF),
    "table_row1": RGBColor(0xFF, 0xFF, 0xFF),
    "table_row2": RGBColor(0xF5, 0xF5, 0xF5),
    "table_border": RGBColor(0xDD, 0xDD, 0xDD),
}

# Matplotlib palette — muted corporate tones with gold lead
MPL_COLORS = ["#C2A269", "#2D2E31", "#5B8C5A", "#4A7B9D", "#C66228",
              "#7B6B8D", "#3B7A57", "#8B6F47", "#526B7A", "#9E7C5B"]

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color):
    """Set a solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=11,
                color=None, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or COLORS["text_dark"]
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_section_header(slide, number, title, top=Inches(0.5)):
    """Add a gold section header matching the Inspired Testing style."""
    # Section number in gold
    if number:
        add_textbox(slide, Inches(0.7), top, Inches(0.6), Inches(0.5),
                    str(number), font_size=18, bold=False,
                    color=COLORS["gold"], font_name="Helvetica")
    # Section title in gold
    x = Inches(1.1) if number else Inches(0.7)
    add_textbox(slide, x, top, Inches(10), Inches(0.5),
                title.upper(), font_size=18, bold=True,
                color=COLORS["gold"], font_name="Helvetica")

    # Gold underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.7), top + Inches(0.55),
                                   Inches(11.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["gold"]
    line.line.fill.background()


def add_sub_heading(slide, number, title, left, top):
    """Add a charcoal sub-heading (e.g. '1.1 HIGH-LEVEL FINDINGS')."""
    text = f"{number}  {title.upper()}" if number else title.upper()
    add_textbox(slide, left, top, Inches(10), Inches(0.4),
                text, font_size=14, bold=True,
                color=COLORS["charcoal"], font_name="Helvetica")


def add_footer(slide):
    """Add a subtle footer with confidential notice."""
    add_textbox(slide, Inches(4.5), SLIDE_HEIGHT - Inches(0.45),
                Inches(4.5), Inches(0.3),
                "C O N F I D E N T I A L", font_size=8,
                color=COLORS["text_muted"], alignment=PP_ALIGN.CENTER,
                font_name="Calibri")


def format_number(n):
    """Format large numbers with K/M suffixes."""
    if not isinstance(n, (int, float)):
        try:
            n = float(n)
        except (ValueError, TypeError):
            return str(n)
    if n >= 1_000_000:
        return f"{n/1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n/1_000:.1f}K"
    return f"{int(n)}"


def create_chart_image(output_path: str, chart_fn, figsize=(8, 4.5)):
    """Create a chart image with clean corporate styling."""
    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.tick_params(colors="#2D2E31", labelsize=9)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["bottom"].set_color("#CCCCCC")
    ax.spines["left"].set_color("#CCCCCC")
    ax.xaxis.label.set_color("#2D2E31")
    ax.yaxis.label.set_color("#2D2E31")
    ax.title.set_color("#2D2E31")
    ax.grid(axis="y", color="#EEEEEE", linewidth=0.5)

    chart_fn(fig, ax)

    plt.tight_layout()
    fig.savefig(output_path, dpi=150, facecolor="white",
                bbox_inches="tight", pad_inches=0.3)
    plt.close(fig)


def add_styled_table(slide, rows, col_widths, left, top, font_size=10):
    """Add a formatted table with corporate styling."""
    n_rows = len(rows)
    n_cols = len(rows[0]) if rows else 0
    if n_rows == 0 or n_cols == 0:
        return None

    total_width = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                          Emu(total_width), Inches(0.3 * n_rows))
    table = table_shape.table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Emu(w)

    for ri, row_data in enumerate(rows):
        is_header = ri == 0
        if is_header:
            bg_color = COLORS["table_header"]
            text_color = COLORS["table_header_text"]
        else:
            bg_color = COLORS["table_row1"] if ri % 2 == 1 else COLORS["table_row2"]
            text_color = COLORS["text_dark"]

        for ci, cell_val in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = str(cell_val)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size if not is_header else font_size + 1)
                paragraph.font.color.rgb = text_color
                paragraph.font.bold = is_header
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.LEFT if ci <= 1 else PP_ALIGN.RIGHT

            cf = cell.fill
            cf.solid()
            cf.fore_color.rgb = bg_color

    return table_shape


# ---------------------------------------------------------------------------
# Individual slide builders
# ---------------------------------------------------------------------------

def build_title_slide(prs, analysis, logo_path=None):
    """Slide 1: Title slide — corporate with gold accent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, COLORS["bg_white"])

    # Gold accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["gold"]
    bar.line.fill.background()

    # Inspired Testing logo — top-right
    if logo_path and os.path.exists(logo_path):
        logo_h = Inches(0.9)
        logo_w = Inches(1.6)  # ~630:354 aspect ratio
        slide.shapes.add_picture(logo_path,
                                  SLIDE_WIDTH - logo_w - Inches(0.5),
                                  Inches(0.2),
                                  logo_w, logo_h)

    # Dark band in upper portion
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   0, Inches(2.0), SLIDE_WIDTH, Inches(3.2))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["charcoal"]
    band.line.fill.background()

    # Title in gold over dark band
    title = os.getenv("REPORT_TITLE", "AI YouTube Trend Analysis")
    add_textbox(slide, Inches(1.5), Inches(2.4), Inches(10), Inches(1.0),
                title.upper(), font_size=36, bold=True,
                alignment=PP_ALIGN.CENTER, color=COLORS["gold"],
                font_name="Helvetica")

    # Subtitle
    period = analysis.get("data_period", {})
    date_text = f"{period.get('from', '?')}  to  {period.get('to', '?')}"
    add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.5),
                date_text, font_size=18, alignment=PP_ALIGN.CENTER,
                color=COLORS["white"], font_name="Helvetica")

    # Stats line
    stats_text = (f"{analysis.get('total_videos_analyzed', 0)} videos analyzed  |  "
                  f"{analysis.get('total_channels_analyzed', 0)} channels  |  "
                  f"Powered by YouTube Data API v3")
    add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
                stats_text, font_size=12, alignment=PP_ALIGN.CENTER,
                color=COLORS["text_muted"])

    # Bottom gold bar
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   0, SLIDE_HEIGHT - Inches(0.06),
                                   SLIDE_WIDTH, Inches(0.06))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = COLORS["gold"]
    bar2.line.fill.background()


def build_summary_slide(prs, analysis):
    """Slide 2: Executive Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_white"])

    add_section_header(slide, "1", "Executive Summary")

    summary = analysis.get("executive_summary", [])
    for i, bullet in enumerate(summary[:7]):
        y = Inches(1.3) + Inches(i * 0.75)

        # Gold bullet marker
        add_textbox(slide, Inches(0.9), y, Inches(0.4), Inches(0.4),
                    "\u25B8", font_size=14, color=COLORS["gold"],
                    font_name="Helvetica")

        add_textbox(slide, Inches(1.3), y, Inches(10.5), Inches(0.6),
                    bullet, font_size=13, color=COLORS["text_dark"])

    add_footer(slide)


def build_top_videos_slide(prs, analysis):
    """Slide 3: Top Trending Videos table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_white"])

    add_section_header(slide, "2", "Top Trending Videos")

    videos = analysis.get("top_videos", {}).get("by_views", [])[:10]
    if not videos:
        add_textbox(slide, Inches(1), Inches(2), Inches(8), Inches(1),
                    "No video data available", font_size=14, color=COLORS["text_muted"])
        return

    rows = [["#", "Title", "Channel", "Views", "Engagement", "Views/Day"]]
    for i, v in enumerate(videos):
        title = str(v.get("title", ""))[:55]
        if len(str(v.get("title", ""))) > 55:
            title += "..."
        rows.append([
            str(i + 1),
            title,
            str(v.get("channel_title", ""))[:25],
            format_number(v.get("view_count", 0)),
            f"{float(v.get('engagement_ratio', 0)):.2%}",
            format_number(v.get("views_per_day", 0)),
        ])

    col_w = [Inches(0.4), Inches(5.5), Inches(2.5), Inches(1.3), Inches(1.3), Inches(1.3)]
    add_styled_table(slide, rows, [Emu(int(w)) for w in col_w],
                     Inches(0.5), Inches(1.3), font_size=10)
    add_footer(slide)


def build_top_channels_slide(prs, analysis):
    """Slide 4: Top Channels table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_white"])

    add_section_header(slide, "3", "Top Channels")

    channels = analysis.get("top_channels", [])[:10]
    if not channels:
        add_textbox(slide, Inches(1), Inches(2), Inches(8), Inches(1),
                    "No channel data available", font_size=14, color=COLORS["text_muted"])
        return

    rows = [["#", "Channel", "Subscribers", "Videos (Dataset)", "Avg Engagement", "Avg Views"]]
    for i, ch in enumerate(channels):
        rows.append([
            str(i + 1),
            str(ch.get("channel_title", ""))[:30],
            format_number(ch.get("subscriber_count", 0)),
            str(int(ch.get("videos_in_dataset", 0))),
            f"{float(ch.get('avg_engagement', 0)):.2%}",
            format_number(ch.get("avg_views", 0)),
        ])

    col_w = [Inches(0.4), Inches(3.5), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)]
    add_styled_table(slide, rows, [Emu(int(w)) for w in col_w],
                     Inches(0.5), Inches(1.3), font_size=10)
    add_footer(slide)


def build_engagement_chart_slide(prs, analysis, tmp_dir):
    """Slide 5: Engagement Analysis bar chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_white"])

    add_section_header(slide, "4", "Engagement Analysis")

    videos = analysis.get("top_videos", {}).get("by_engagement", [])[:15]
    if not videos:
        return

    chart_path = os.path.join(tmp_dir, "chart_engagement.png")

    def draw(fig, ax):
        titles = [v.get("title", "")[:40] + ("..." if len(v.get("title", "")) > 40 else "")
                  for v in reversed(videos)]
        values = [float(v.get("engagement_ratio", 0)) * 100 for v in reversed(videos)]

        bars = ax.barh(range(len(titles)), values, color="#C2A269", height=0.7,
                       edgecolor="white", linewidth=0.5)
        ax.set_yticks(range(len(titles)))
        ax.set_yticklabels(titles, fontsize=8, color="#2D2E31")
        ax.set_xlabel("Engagement Ratio (%)", fontsize=10)
        ax.set_title("Top 15 Videos by Engagement", fontsize=13,
                     fontweight="bold", pad=12)
        ax.xaxis.set_major_formatter(mticker.FormatStrFormatter("%.1f%%"))

    create_chart_image(chart_path, draw, figsize=(10, 6))
    slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))
    add_footer(slide)


def build_topic_distribution_slide(prs, analysis, tmp_dir):
    """Slide 6: Topic Distribution pie chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_white"])

    add_section_header(slide, "5", "Topic Distribution")

    categories = analysis.get("topic_distribution", {}).get("categories", {})
    if not categories:
        return

    chart_path = os.path.join(tmp_dir, "chart_topics.png")
    sorted_cats = sorted(categories.items(), key=lambda x: x[1], reverse=True)
    labels = [c[0] for c in sorted_cats if c[1] > 0]
    sizes = [c[1] for c in sorted_cats if c[1] > 0]
    colors = [MPL_COLORS[i % len(MPL_COLORS)] for i in range(len(labels))]

    def draw(fig, ax):
        explode = [0.03 if i == 0 else 0 for i in range(len(labels))]
        wedges, texts, autotexts = ax.pie(
            sizes, labels=labels, autopct="%1.0f%%", startangle=140,
            colors=colors, explode=explode, pctdistance=0.80,
            textprops={"fontsize": 9, "color": "#2D2E31"},
        )
        for t in autotexts:
            t.set_color("#2D2E31")
            t.set_fontsize(8)
            t.set_fontweight("bold")
        ax.set_title("Videos by Topic Category", fontsize=13,
                     fontweight="bold", color="#2D2E31", pad=15)

    create_chart_image(chart_path, draw, figsize=(7, 5))
    slide.shapes.add_picture(chart_path, Inches(0.3), Inches(1.1), Inches(7), Inches(5.8))

    # Engagement table on the right
    eng = analysis.get("topic_distribution", {}).get("engagement_by_category", {})
    if eng:
        add_sub_heading(slide, "5.1", "Engagement by Topic", Inches(8), Inches(1.3))
        rows = [["Topic", "Engagement"]]
        for cat in labels[:8]:
            val = eng.get(cat, 0)
            rows.append([cat, f"{val:.2%}"])

        col_w = [Inches(2.8), Inches(1.5)]
        add_styled_table(slide, rows, [Emu(int(w)) for w in col_w],
                         Inches(8), Inches(1.8), font_size=11)

    add_footer(slide)


def build_velocity_slide(prs, analysis, tmp_dir):
    """Slide 7: View Velocity Trends line chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_white"])

    add_section_header(slide, "6", "View Velocity Trends")

    weekly = analysis.get("view_velocity", {}).get("weekly_data", [])
    if not weekly:
        add_textbox(slide, Inches(1), Inches(2), Inches(8), Inches(1),
                    "Insufficient data for velocity analysis", font_size=14,
                    color=COLORS["text_muted"])
        return

    chart_path = os.path.join(tmp_dir, "chart_velocity.png")

    def draw(fig, ax):
        weeks = [w["week_label"] for w in weekly]
        vpd = [w["avg_views_per_day"] for w in weekly]

        ax.plot(weeks, vpd, color="#C2A269", linewidth=2.5, marker="o",
                markersize=7, markerfacecolor="#C2A269", markeredgecolor="#2D2E31",
                markeredgewidth=1)
        ax.fill_between(range(len(weeks)), vpd, alpha=0.10, color="#C2A269")
        ax.set_xlabel("Week", fontsize=10)
        ax.set_ylabel("Avg Views per Day", fontsize=10)
        ax.set_title("Weekly Average View Velocity", fontsize=13,
                     fontweight="bold", pad=12)
        plt.xticks(rotation=45, ha="right", fontsize=8)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format_number(x)))

    create_chart_image(chart_path, draw, figsize=(10, 5))
    slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))
    add_footer(slide)


def build_patterns_slide(prs, analysis, tmp_dir):
    """Slide 8: Content Patterns — 4 mini-charts."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_white"])

    add_section_header(slide, "7", "Content Patterns")

    patterns = analysis.get("content_patterns", {})

    # --- Duration chart (top-left) ---
    dur = patterns.get("duration_performance", {})
    if dur:
        dur_path = os.path.join(tmp_dir, "chart_duration.png")

        def draw_dur(fig, ax):
            labels_d = list(dur.keys())
            views_d = [dur[k].get("avg_views", 0) for k in labels_d]
            ax.bar(labels_d, views_d, color=MPL_COLORS[:len(labels_d)], width=0.6,
                   edgecolor="white", linewidth=0.5)
            ax.set_title("Avg Views by Duration", fontsize=11, fontweight="bold")
            ax.set_ylabel("Avg Views", fontsize=9)
            ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format_number(x)))
            plt.xticks(fontsize=8, rotation=20)

        create_chart_image(dur_path, draw_dur, figsize=(5, 3.2))
        slide.shapes.add_picture(dur_path, Inches(0.3), Inches(1.2), Inches(6), Inches(3))

    # --- Publish day chart (top-right) ---
    day_perf = patterns.get("publish_day_performance", {})
    if day_perf:
        day_path = os.path.join(tmp_dir, "chart_publish_day.png")

        def draw_day(fig, ax):
            days = list(day_perf.keys())
            views_d = [day_perf[d].get("avg_views", 0) for d in days]
            ax.bar(days, views_d, color="#C2A269", width=0.6,
                   edgecolor="white", linewidth=0.5)
            ax.set_title("Avg Views by Publish Day", fontsize=11, fontweight="bold")
            ax.set_ylabel("Avg Views", fontsize=9)
            ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format_number(x)))
            plt.xticks(fontsize=7, rotation=30)

        create_chart_image(day_path, draw_day, figsize=(5, 3.2))
        slide.shapes.add_picture(day_path, Inches(6.7), Inches(1.2), Inches(6), Inches(3))

    # --- Top keywords (bottom-left) ---
    keywords = patterns.get("top_title_words_top25pct", [])
    if keywords:
        kw_path = os.path.join(tmp_dir, "chart_keywords.png")

        def draw_kw(fig, ax):
            words = [w["word"] for w in keywords[:12]]
            counts = [w["count"] for w in keywords[:12]]
            ax.barh(range(len(words)), counts, color=MPL_COLORS[:len(words)],
                    height=0.6, edgecolor="white", linewidth=0.5)
            ax.set_yticks(range(len(words)))
            ax.set_yticklabels(words, fontsize=9)
            ax.set_title("Top Keywords in High-Performing Titles",
                         fontsize=11, fontweight="bold")
            ax.invert_yaxis()

        create_chart_image(kw_path, draw_kw, figsize=(5, 3.2))
        slide.shapes.add_picture(kw_path, Inches(0.3), Inches(4.3), Inches(6), Inches(3))

    # --- Title length insight (bottom-right) ---
    title_data = patterns.get("title_length", {})
    num_data = patterns.get("number_in_title", {})
    if title_data or num_data:
        # Light card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(7), Inches(4.4), Inches(5.5), Inches(2.6))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["bg_light"]
        card.line.color.rgb = COLORS["table_border"]
        card.line.width = Pt(0.5)

        y_pos = Inches(4.6)
        add_sub_heading(slide, "7.1", "Title Insights", Inches(7.3), y_pos)

        if title_data:
            add_textbox(slide, Inches(7.3), y_pos + Inches(0.5), Inches(5), Inches(0.4),
                        f"Avg title length (all): {title_data.get('avg_title_length_all', 0):.0f} chars",
                        font_size=12, color=COLORS["text_dark"])
            add_textbox(slide, Inches(7.3), y_pos + Inches(0.9), Inches(5), Inches(0.4),
                        f"Avg title length (top 25%): {title_data.get('avg_title_length_top25pct', 0):.0f} chars",
                        font_size=12, color=COLORS["gold"], bold=True)

        if num_data:
            add_textbox(slide, Inches(7.3), y_pos + Inches(1.4), Inches(5), Inches(0.4),
                        f"Titles with numbers: {num_data.get('pct_with_number', 0):.0f}%",
                        font_size=12, color=COLORS["text_dark"])
            views_with = format_number(num_data.get("avg_views_with", 0))
            views_without = format_number(num_data.get("avg_views_without", 0))
            add_textbox(slide, Inches(7.3), y_pos + Inches(1.8), Inches(5), Inches(0.4),
                        f"Avg views: {views_with} (with #) vs {views_without} (without)",
                        font_size=12, color=COLORS["text_body"])

    add_footer(slide)


def build_trending_slide(prs, analysis):
    """Slide 9: Trending vs Declining Topics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_white"])

    add_section_header(slide, "8", "Trending vs Declining Topics")

    trends = analysis.get("trending_topics", {})
    trending = trends.get("trending", [])
    declining = trends.get("declining", [])

    # Left column: Trending Up
    card_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = COLORS["bg_light"]
    card_l.line.color.rgb = COLORS["table_border"]
    card_l.line.width = Pt(0.5)

    add_sub_heading(slide, "", "Trending Up", Inches(0.8), Inches(1.5))

    for i, t in enumerate(trending[:6]):
        y = Inches(2.1) + Inches(i * 0.75)
        add_textbox(slide, Inches(1.0), y, Inches(0.5), Inches(0.4),
                    "\u25B2", font_size=14, color=COLORS["success"])
        add_textbox(slide, Inches(1.5), y, Inches(3.5), Inches(0.4),
                    t["topic"], font_size=13, color=COLORS["text_dark"])
        add_textbox(slide, Inches(5.0), y, Inches(1.2), Inches(0.4),
                    f"+{t['change_pct']}%", font_size=13, bold=True,
                    color=COLORS["success"], alignment=PP_ALIGN.RIGHT)

    if not trending:
        add_textbox(slide, Inches(1.0), Inches(2.5), Inches(4), Inches(0.4),
                    "No significant upward trends detected",
                    font_size=12, color=COLORS["text_muted"])

    # Right column: Declining
    card_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLORS["bg_light"]
    card_r.line.color.rgb = COLORS["table_border"]
    card_r.line.width = Pt(0.5)

    add_sub_heading(slide, "", "Declining", Inches(7.3), Inches(1.5))

    for i, t in enumerate(declining[:6]):
        y = Inches(2.1) + Inches(i * 0.75)
        add_textbox(slide, Inches(7.5), y, Inches(0.5), Inches(0.4),
                    "\u25BC", font_size=14, color=COLORS["warning"])
        add_textbox(slide, Inches(8.0), y, Inches(3.5), Inches(0.4),
                    t["topic"], font_size=13, color=COLORS["text_dark"])
        add_textbox(slide, Inches(11.5), y, Inches(1.2), Inches(0.4),
                    f"{t['change_pct']}%", font_size=13, bold=True,
                    color=COLORS["warning"], alignment=PP_ALIGN.RIGHT)

    if not declining:
        add_textbox(slide, Inches(7.5), Inches(2.5), Inches(4), Inches(0.4),
                    "No significant downward trends detected",
                    font_size=12, color=COLORS["text_muted"])

    add_footer(slide)


def build_recommendations_slide(prs, analysis):
    """Slide 10: Recommendations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_white"])

    add_section_header(slide, "9", "Recommendations")

    recs = analysis.get("recommendations", [])[:6]

    for i, rec in enumerate(recs):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(1.3) + row * Inches(2.0)

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, y, Inches(5.9), Inches(1.7))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["bg_light"]
        card.line.color.rgb = COLORS["table_border"]
        card.line.width = Pt(0.5)

        # Gold accent bar on left
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x + Inches(0.05), y + Inches(0.15),
                                      Inches(0.06), Inches(1.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS["gold"]
        bar.line.fill.background()

        # Number in gold
        add_textbox(slide, x + Inches(0.25), y + Inches(0.15), Inches(0.5), Inches(0.4),
                    str(i + 1), font_size=20, bold=True, color=COLORS["gold"],
                    font_name="Helvetica")

        # Recommendation text
        add_textbox(slide, x + Inches(0.75), y + Inches(0.15), Inches(4.9), Inches(1.4),
                    rec, font_size=11, color=COLORS["text_dark"])

    add_footer(slide)


def build_appendix_slide(prs, analysis):
    """Slide 11: Appendix — Methodology."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["bg_white"])

    add_section_header(slide, "10", "Appendix — Methodology")

    period = analysis.get("data_period", {})
    info_lines = [
        ("Data Source:", "YouTube Data API v3"),
        ("Date Range:", f"{period.get('from', '?')} to {period.get('to', '?')}"),
        ("Total Videos Analyzed:", str(analysis.get("total_videos_analyzed", 0))),
        ("Total Channels Analyzed:", str(analysis.get("total_channels_analyzed", 0))),
        ("Generated:", analysis.get("generated_at", "N/A")[:19]),
    ]

    for i, (label, value) in enumerate(info_lines):
        y = Inches(1.4) + Inches(i * 0.45)
        add_textbox(slide, Inches(0.9), y, Inches(3), Inches(0.4),
                    label, font_size=12, bold=True, color=COLORS["charcoal"])
        add_textbox(slide, Inches(4.0), y, Inches(7), Inches(0.4),
                    value, font_size=12, color=COLORS["text_dark"])

    # Formulas section
    add_sub_heading(slide, "10.1", "Metrics Definitions", Inches(0.7), Inches(4.0))

    formulas = [
        "Engagement Ratio = (Likes + Comments) / Views",
        "View Velocity = Total Views / Days Since Published",
        "Topic classification uses keyword matching on video titles and tags.",
    ]
    for i, formula in enumerate(formulas):
        y = Inches(4.5) + Inches(i * 0.4)
        add_textbox(slide, Inches(0.9), y, Inches(10), Inches(0.35),
                    formula, font_size=11, color=COLORS["text_body"])

    # Disclaimer
    add_textbox(slide, Inches(0.9), Inches(6.0), Inches(10), Inches(0.5),
                "Disclaimer: Data reflects publicly available YouTube metrics at the "
                "time of collection. Metrics may change as videos continue to accrue views.",
                font_size=10, color=COLORS["text_muted"])

    add_footer(slide)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main(data_dir: str = None) -> dict:
    """
    Generate a PowerPoint slide deck from analysis results.

    Args:
        data_dir: Directory containing analysis data. Defaults to .tmp/.

    Returns:
        dict with status and path to generated .pptx file.
    """
    load_dotenv()

    if data_dir is None:
        data_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), ".tmp")

    manifest_path = os.path.join(data_dir, "youtube_latest.json")

    try:
        if not os.path.exists(manifest_path):
            return {"status": "error",
                    "error": f"Manifest not found at {manifest_path}. Run youtube_analyzer.py first."}

        with open(manifest_path, "r", encoding="utf-8") as f:
            manifest = json.load(f)

        analysis_file = manifest.get("latest_analysis", {}).get("analysis_file")
        if not analysis_file or not os.path.exists(analysis_file):
            return {"status": "error",
                    "error": f"Analysis file not found: {analysis_file}. Run youtube_analyzer.py first."}

        with open(analysis_file, "r", encoding="utf-8") as f:
            analysis = json.load(f)

        print("  Creating presentation ...")

        # Create widescreen presentation (13.333 x 7.5 inches)
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        # Build slides
        # Resolve logo path
        logo_path = os.path.join(data_dir, "inspired_testing_logo.png")
        if not os.path.exists(logo_path):
            logo_path = None

        print("    Slide 1: Title ...")
        build_title_slide(prs, analysis, logo_path=logo_path)

        print("    Slide 2: Executive Summary ...")
        build_summary_slide(prs, analysis)

        print("    Slide 3: Top Videos ...")
        build_top_videos_slide(prs, analysis)

        print("    Slide 4: Top Channels ...")
        build_top_channels_slide(prs, analysis)

        print("    Slide 5: Engagement Analysis ...")
        build_engagement_chart_slide(prs, analysis, data_dir)

        print("    Slide 6: Topic Distribution ...")
        build_topic_distribution_slide(prs, analysis, data_dir)

        print("    Slide 7: View Velocity ...")
        build_velocity_slide(prs, analysis, data_dir)

        print("    Slide 8: Content Patterns ...")
        build_patterns_slide(prs, analysis, data_dir)

        print("    Slide 9: Trending vs Declining ...")
        build_trending_slide(prs, analysis)

        print("    Slide 10: Recommendations ...")
        build_recommendations_slide(prs, analysis)

        print("    Slide 11: Appendix ...")
        build_appendix_slide(prs, analysis)

        # Save
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        report_path = os.path.join(data_dir, f"youtube_trend_report_{ts}.pptx")
        prs.save(report_path)

        file_size = os.path.getsize(report_path)
        print(f"  Saved: {report_path} ({file_size / 1024:.0f} KB)")

        # Update manifest
        manifest["latest_report"] = {"report_file": report_path, "timestamp": ts}
        with open(manifest_path, "w", encoding="utf-8") as f:
            json.dump(manifest, f, indent=2)

        return {
            "status": "success",
            "data": {
                "report_path": report_path,
                "slide_count": 11,
                "file_size_kb": round(file_size / 1024, 1),
            },
        }

    except Exception as e:
        return {"status": "error", "error": str(e)}


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="YouTube Trend Report Slide Generator")
    parser.add_argument("--data-dir", type=str, help="Directory containing analysis data")
    args = parser.parse_args()

    result = main(data_dir=args.data_dir)

    if result["status"] == "success":
        print(f"\nSuccess: Generated {result['data']['slide_count']}-slide deck")
        print(f"  File: {result['data']['report_path']}")
        print(f"  Size: {result['data']['file_size_kb']} KB")
        sys.exit(0)
    else:
        print(f"\nError: {result['error']}", file=sys.stderr)
        sys.exit(1)
